import logging
import time
from fastapi import APIRouter, File, UploadFile, HTTPException, Depends, BackgroundTasks
from typing import List
from .schemas import QueryRequest, QueryResponse, StatusResponse
from .auth import get_current_user
from storage.ingestion import ingestion_pipeline
from storage.vector_db import vector_db
from storage.memory_cache import memory_cache
from storage.graph_db import graph_db
from agents.orchestrator import process_query_workflow
from core.embeddings import LocalEmbedder
from immune.macrophage import immune_system
from core.confluence import confluence_sync

logger = logging.getLogger(__name__)
api_router = APIRouter()

def process_upload_in_background(filename: str, text_content: str, user_id: int):
    try:
        nodes = ingestion_pipeline.process_document(filename, text_content)
        embedder = LocalEmbedder.get_embedder()
        texts = [n["text"] for n in nodes]
        real_vectors = embedder.embed_documents(texts)
        payloads = [
            {"text": n["text"], "file_name": filename, "user_id": user_id} 
            for n in nodes
        ]
        vector_db.upsert_vectors(real_vectors, payloads)
        # Populate graph DB with keyword relationships from this document
        ingestion_pipeline.populate_graph(filename, nodes)
        logger.info(f"Background processing complete for {filename}: {len(nodes)} chunks indexed.")
    except Exception as e:
        logger.error(f"Background upload processing failed for {filename}: {e}")

@api_router.post("/upload", response_model=StatusResponse)
async def upload_document(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    current_user: dict = Depends(get_current_user)
):
    """
    Ingests PDF, TXT, DOCX, PPTX with proper text extraction per file type.
    """
    try:
        content = await file.read()
        filename = file.filename
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else "txt"

        # --- Proper text extraction per file format ---
        text_content = ""
        if ext == "pptx":
            try:
                import io
                from pptx import Presentation
                prs = Presentation(io.BytesIO(content))
                slide_texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_texts.append(shape.text.strip())
                text_content = "\n\n".join(slide_texts)
                logger.info(f"Extracted {len(slide_texts)} text blocks from PPTX '{filename}'")
            except Exception as e:
                logger.error(f"PPTX extraction failed: {e}")
                raise HTTPException(status_code=422, detail=f"Could not read PPTX file: {e}")

        elif ext == "pdf":
            try:
                import io
                from pypdf import PdfReader
                reader = PdfReader(io.BytesIO(content))
                pages = [page.extract_text() or "" for page in reader.pages]
                text_content = "\n\n".join(p for p in pages if p.strip())
                logger.info(f"Extracted {len(reader.pages)} pages from PDF '{filename}'")
            except Exception as e:
                logger.error(f"PDF extraction failed: {e}")
                raise HTTPException(status_code=422, detail=f"Could not read PDF file: {e}")

        elif ext == "docx":
            try:
                import io, docx
                doc = docx.Document(io.BytesIO(content))
                text_content = "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
                logger.info(f"Extracted {len(doc.paragraphs)} paragraphs from DOCX '{filename}'")
            except Exception as e:
                logger.warning(f"DOCX extraction failed, falling back to UTF-8: {e}")
                text_content = content.decode("utf-8", errors="ignore")

        else:  # txt, md, csv, etc.
            text_content = content.decode("utf-8", errors="ignore")

        if not text_content.strip():
            raise HTTPException(status_code=422, detail="No readable text could be extracted from this file.")

        file_size = len(text_content)
        logger.info(f"Received file: {filename} (Extracted {file_size} chars) from {current_user['username']}")

        # Small file: store in memory cache directly (bypass RAG)
        if file_size < 15000:
            memory_cache.store_file(filename, text_content)
            background_tasks.add_task(immune_system.scan_for_conflicts, text_content, filename)
            return StatusResponse(
                status="success",
                message=f"Document '{filename}' saved. Macrophage background scan initiated.",
                data={"chunks_created": 0, "status": "direct_injection"}
            )

        # Large file: embed in background
        background_tasks.add_task(process_upload_in_background, filename, text_content, current_user['id'])
        return StatusResponse(
            status="success",
            message=f"Document '{filename}' successfully ingested. Embedding processing running in background.",
            data={"status": "processing"}
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error uploading file: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.post("/query", response_model=QueryResponse)
async def query_knowledge(
    request: QueryRequest,
    current_user: dict = Depends(get_current_user)
):
    """
    Handles natural language queries and triggers the Multi-Agent LangGraph workflow.
    """
    try:
        logger.info(f"Received query: {request.query} in mode: {request.mode} for user: {current_user['username']}")
        
        # Trigger Multi-Agent LangGraph Workflow directly
        start_time = time.time()
        agent_raw_result = process_query_workflow(request.query, mode=request.mode)
        end_time = time.time()
        
        logger.info(f"LangGraph execution completed in {round(end_time-start_time, 2)}s.")
        
        # Extract metadata from state
        sources = list(set(agent_raw_result.get("vector_context", [])))
        if not sources or sources == ["No semantic data found in Qdrant."]:
            sources = ["Knowledge Graph"]
            
        return QueryResponse(
            answer=agent_raw_result.get("final_answer", "No answer could be generated."),
            confidence_score=agent_raw_result.get("confidence", 0.0),
            confidence_level="High" if agent_raw_result.get("confidence", 0) > 85 else "Medium",
            strategy=agent_raw_result.get("strategy", "Strict Hybrid Validation"),
            sources=sources,
            steps=agent_raw_result.get("steps_taken", [])
        )
    except Exception as e:
        logger.error(f"Error executing query: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.delete("/documents/{filename}", response_model=StatusResponse)
async def delete_document(
    filename: str,
    current_user: dict = Depends(get_current_user)
):
    """
    Deletes a document completely from VectorDB and MemoryCache.
    """
    try:
        logger.info(f"User {current_user['username']} requested deletion of {filename}")
        
        # 1. Delete from Vector DB
        vector_db.delete_by_filename(filename)
        
        # 2. Delete from Memory Cache
        memory_cache.delete_file(filename)
        
        return StatusResponse(
            status="success",
            message=f"Document '{filename}' successfully deleted from databases.",
            data={"deleted_file": filename}
        )
    except Exception as e:
        logger.error(f"Error deleting file {filename}: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.get("/immune/status")
async def get_immune_status(current_user: dict = Depends(get_current_user)):
    """Exposes current macrophage threat detection state to the frontend"""
    conflicts = immune_system.conflicts
    return {
        "status": "infected" if len(conflicts) > 0 else "secure",
        "active_threats": len(conflicts),
        "reports": conflicts
    }

@api_router.get("/suggestions")
async def get_dynamic_suggestions(current_user: dict = Depends(get_current_user)):
    """Fetches random graph entities and dynamically templates them into suggested questions"""
    entities = graph_db.get_random_entities(5)
    
    if not entities:
        return {"suggestions": [
            "What documents are currently uploaded?",
            "Can you summarize the main concepts?",
            "What security policies are available?"
        ]}
        
    import random
    templates = [
        "Explain the significance of '{0}' in the documents.",
        "What are the main findings regarding '{0}'?",
        "Can you summarize the context around '{0}'?",
        "How is '{0}' related to other concepts?"
    ]
    
    suggestions = [tmpl.format(entity) for entity, tmpl in zip(chosen, [random.choice(templates) for _ in chosen])]
    return {"suggestions": suggestions}

@api_router.post("/confluence/sync", response_model=StatusResponse)
async def sync_confluence(
    request: dict, # Basic dict for URL input
    current_user: dict = Depends(get_current_user)
):
    url = request.get("url")
    if not url:
        raise HTTPException(status_code=400, detail="Confluence URL is required")
    
    # Run sync
    result = await confluence_sync.sync_page(url, current_user['id'])
    if result["status"] == "error":
        raise HTTPException(status_code=500, detail=result["message"])
    
    return StatusResponse(status="success", message=f"Ingested {result['title']} ({result['chars']} chars)")

@api_router.get("/graph/stats")
async def get_graph_stats(current_user: dict = Depends(get_current_user)):
    """Debug endpoint to inspect the knowledge graph DB state."""
    nodes = list(graph_db.graph.nodes())
    edges = list(graph_db.graph.edges(data=True))
    sample_nodes = nodes[:20] if nodes else []
    sample_edges = [{"from": u, "to": v, "rel": d.get("relation", "?")} for u, v, d in edges[:20]]
    return {
        "total_nodes": len(nodes),
        "total_edges": len(edges),
        "sample_nodes": sample_nodes,
        "sample_edges": sample_edges,
        "memory_cache_files": list(memory_cache.cache.keys()),
        "memory_cache_total_chars": sum(len(v) for v in memory_cache.cache.values()),
    }
